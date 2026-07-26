import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "report", "presentation.pptx")

# Charte reprise du template fourni (Work-Life Integration, Slidesgo) : orange
# terracotta, encre, papier, blanc. Polices de substitution universellement
# disponibles, dans le même esprit que le template (sans-serif gras pour les
# titres, serif italique pour le mot-accent).
ORANGE = RGBColor(0xB4, 0x3B, 0x13)
INK = RGBColor(0x25, 0x25, 0x25)
PAPER = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x8A, 0x7E, 0x74)

FONT_HEAD = "Arial"
FONT_ACCENT = "Georgia"
FONT_BODY = "Arial"

GROUP = "Mohamed Elyes Bahouri · Mathis Te · Olade Mac-Aurel Laourou Babalekon"
FOOTER_LEFT = "Reinforcement Learning · DP, Monte Carlo, TD, Planning"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

page_counter = 0


def new_slide(bg):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_rule(slide, left, top, width, color=ORANGE, weight=1.25):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # Retire le style de thème par défaut (ombre portée) pour un filet net et plat.
    style_el = line._element.find(qn("p:style"))
    if style_el is not None:
        line._element.remove(style_el)
    return line


def add_text(slide, left, top, width, height, text, size=16, color=INK, bold=False,
             italic=False, font=FONT_BODY, align=PP_ALIGN.LEFT, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_headline(slide, left, top, width, height, parts, size=34, color=WHITE,
                  align=PP_ALIGN.LEFT, line_spacing=1.05):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first_line = True
    for line_parts in parts:
        p = tf.paragraphs[0] if first_line else tf.add_paragraph()
        first_line = False
        p.alignment = align
        p.line_spacing = line_spacing
        for text, italic in line_parts:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = not italic
            run.font.italic = italic
            run.font.name = FONT_ACCENT if italic else FONT_HEAD
            run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=14, color=INK,
                 marker_color=ORANGE, space_after=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.color.rgb = marker_color
        r1.font.bold = True
        r1.font.size = Pt(size)
        r1.font.name = FONT_BODY
        r2 = p.add_run()
        r2.text = item
        r2.font.color.rgb = color
        r2.font.size = Pt(size)
        r2.font.name = FONT_BODY
    return box


def add_footer(slide, dark_bg=False):
    global page_counter
    page_counter += 1
    color = WHITE if dark_bg else MUTED
    line_color = WHITE if dark_bg else ORANGE
    add_rule(slide, Inches(0.55), Inches(5.18), Inches(8.9), color=line_color, weight=0.75)
    add_text(slide, Inches(0.55), Inches(5.24), Inches(6.5), Inches(0.3), FOOTER_LEFT,
              size=8, color=color, font=FONT_BODY)
    add_text(slide, Inches(8.6), Inches(5.24), Inches(0.85), Inches(0.3), f"{page_counter:02d}",
              size=8, color=color, font=FONT_BODY, align=PP_ALIGN.RIGHT)


def add_kicker(slide, text, color=ORANGE, left=Inches(0.55), top=Inches(0.4)):
    add_text(slide, left, top, Inches(8.5), Inches(0.3), text, size=11, color=color,
              bold=True, font=FONT_BODY)


def add_table(slide, left, top, width, height, header, rows, col_widths=None,
              body_size=12, header_size=12):
    n_rows = len(rows) + 1
    n_cols = len(header)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, text in enumerate(header):
        cell = table.cell(0, j)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_BODY
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    for i, row in enumerate(rows, start=1):
        for j, text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else PAPER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(body_size)
            run.font.color.rgb = INK
            run.font.name = FONT_BODY
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table


def divider_slide(number, title_parts, kicker):
    slide = new_slide(PAPER)
    add_kicker(slide, kicker)
    add_text(slide, Inches(0.55), Inches(1.5), Inches(3), Inches(1.2), number,
              size=64, color=ORANGE, italic=True, font=FONT_ACCENT)
    add_headline(slide, Inches(0.55), Inches(3.15), Inches(8.9), Inches(1.2),
                  title_parts, size=32, color=INK)
    add_rule(slide, Inches(0.55), Inches(4.55), Inches(8.9), color=ORANGE, weight=1.5)
    add_footer(slide, dark_bg=False)
    return slide


# ---------------------------------------------------------------------------
# 1. Titre
# ---------------------------------------------------------------------------
slide = new_slide(ORANGE)
add_kicker(slide, "2026-4A-IABD · (Deep) Reinforcement Learning P1", color=WHITE)
add_headline(
    slide, Inches(0.55), Inches(1.55), Inches(8.9), Inches(2.4),
    [
        [("APPRENDRE À ", False)],
        [("décider", True), (" PAR RENFORCEMENT", False)],
    ],
    size=40, color=WHITE,
)
add_text(
    slide, Inches(0.55), Inches(3.55), Inches(8.9), Inches(0.6),
    "Programmation dynamique, Monte Carlo, Temporal Difference Learning et Planning, "
    "testés sur cinq environnements avec huit algorithmes.",
    size=15, color=WHITE, italic=True, font=FONT_BODY,
)
add_rule(slide, Inches(0.55), Inches(5.0), Inches(8.9), color=WHITE, weight=1)
add_text(slide, Inches(0.55), Inches(5.08), Inches(6.5), Inches(0.35), GROUP, size=10, color=WHITE, font=FONT_BODY)
add_text(slide, Inches(8.0), Inches(5.08), Inches(1.45), Inches(0.35), "Juillet 2026", size=10,
          color=WHITE, font=FONT_BODY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# 2. Objectif & méthodologie
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Introduction")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("OBJECTIF & ", False), ("méthodologie", True)]], size=28, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_bullets(slide, Inches(0.55), Inches(1.65), Inches(8.9), Inches(3.4), [
    "Le projet s'appuie sur deux bibliothèques séparées, une pour les algorithmes et une pour "
    "les environnements, qui communiquent via une interface commune.",
    "Pour comparer équitablement, on entraîne chaque algorithme puis on évalue la politique "
    "obtenue sur 200 épisodes joués en glouton, sans apprentissage : le score est ainsi "
    "comparable d'un algorithme à l'autre.",
    "La programmation dynamique est peu sensible au hasard une fois convergée, donc on ne la "
    "lance qu'une seule fois par environnement.",
    "Les méthodes model-free dépendent beaucoup de l'exploration, donc on les relance sur 5 "
    "graines aléatoires pour mesurer un vrai taux de réussite, pas un coup de chance isolé.",
    "Le seuil de réussite n'est pas fixé au hasard : c'est le score optimal réel de "
    "l'environnement, calculé avec Value Iteration, avec une marge de 0,1.",
    "On a aussi relevé epsilon à 0,3 au lieu de 0,1 pour Sarsa, Q-Learning, Dyna-Q et "
    "l'on-policy MC control, à cause d'un problème d'exploration détaillé plus loin.",
], size=12, space_after=7)
add_footer(slide)

# ---------------------------------------------------------------------------
# 3. Divider : environnements
# ---------------------------------------------------------------------------
divider_slide("01.", [[("LES ", False), ("environnements", True)]], "Plan")

# ---------------------------------------------------------------------------
# 4. Les 4 environnements développés
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Environnements")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("QUATRE ENVIRONNEMENTS ", False), ("développés", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))

cards = [
    ("LINE WORLD", "Ligne de 5 cases, départ au milieu. Sortir à gauche donne -1, à droite +1."),
    ("GRID WORLD", "Grille 5×5, 4 déplacements possibles. L'objectif vaut +1, le piège -1, et "
                    "un mur laisse l'agent sur place."),
    ("TWO ROUND RPS", "2 manches de pierre-feuille-ciseaux. Au round 2, l'adversaire est "
                        "obligé de rejouer le coup de l'agent au round 1."),
    ("MONTY HALL (3 ET 5 PORTES)", "L'agent choisit une porte, une porte perdante est "
                                     "retirée, puis il garde ou change. Le score optimal est "
                                     "de 2/3 avec 3 portes et 4/5 avec 5 portes."),
]
positions = [(0.55, 1.8), (5.05, 1.8), (0.55, 3.55), (5.05, 3.55)]
for (title, desc), (l, t) in zip(cards, positions):
    add_text(slide, Inches(l), Inches(t), Inches(4.3), Inches(0.3), title, size=13.5, bold=True,
              color=ORANGE, font=FONT_BODY)
    add_text(slide, Inches(l), Inches(t + 0.35), Inches(4.3), Inches(1.3), desc, size=11.5,
              color=INK, font=FONT_BODY, line_spacing=1.15)
add_footer(slide)

# ---------------------------------------------------------------------------
# 5. Environnements secrets
# ---------------------------------------------------------------------------
slide = new_slide(PAPER)
add_kicker(slide, "Environnements")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("LES ENVIRONNEMENTS ", False), ("secrets", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_text(slide, Inches(0.55), Inches(1.65), Inches(8.9), Inches(0.8),
          "Ce sont des bibliothèques compilées fournies par le cours (.so, .dll, .dylib), dont "
          "on ne connaît pas les règles internes. Un adaptateur, SecretEnvAdapter, les rend "
          "compatibles avec l'interface commune du projet.",
          size=12.5, color=INK, line_spacing=1.15)
add_table(
    slide, Inches(0.55), Inches(2.55), Inches(8.9), Inches(1.5),
    ["Environnement", "États", "Actions"],
    [
        ["Secret Env 0", "8 192", "3"],
        ["Secret Env 1", "65 536", "3"],
        ["Secret Env 2", "2 097 152", "3"],
        ["Secret Env 3", "65 536", "3"],
    ],
    col_widths=[Inches(4.5), Inches(2.7), Inches(1.7)],
)
add_bullets(slide, Inches(0.55), Inches(4.2), Inches(8.9), Inches(0.9), [
    "On ne peut pas replacer l'agent directement dans un état choisi, donc Monte Carlo ES "
    "ne fonctionne pas ici.",
    "L'espace d'états est trop grand pour notre implémentation de la programmation "
    "dynamique, donc Policy et Value Iteration ne sont pas utilisés non plus.",
], size=11.5)
add_footer(slide)

# ---------------------------------------------------------------------------
# 6. Divider : algorithmes
# ---------------------------------------------------------------------------
divider_slide("02.", [[("LES ", False), ("algorithmes", True)]], "Plan")

# ---------------------------------------------------------------------------
# 7. Les 4 familles d'algorithmes
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Algorithmes")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("QUATRE FAMILLES, ", False), ("huit algorithmes", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_table(
    slide, Inches(0.55), Inches(1.7), Inches(8.9), Inches(3.2),
    ["Famille", "Algorithmes", "Principe"],
    [
        ["Dynamic Programming", "Policy Iteration, Value Iteration",
         "Le modèle est connu à l'avance, donc on calcule la solution sans jouer une seule partie."],
        ["Monte Carlo", "MC ES, on-policy first-visit, off-policy",
         "On joue des épisodes complets et on apprend à partir des retours obtenus à la fin."],
        ["Temporal Difference", "Sarsa, Q-Learning",
         "On corrige Q après chaque action, sans attendre la fin de l'épisode."],
        ["Planning", "Dyna-Q",
         "C'est Q-Learning complété par un modèle appris, rejoué en interne pour apprendre plus vite."],
    ],
    col_widths=[Inches(2.1), Inches(2.5), Inches(4.3)],
    body_size=11.5, header_size=12,
)
add_footer(slide)

# ---------------------------------------------------------------------------
# 8. Divider : résultats
# ---------------------------------------------------------------------------
divider_slide("03.", [[("RÉSULTATS ", False), ("comparés", True)]], "Plan")

# ---------------------------------------------------------------------------
# 9. Résultats sur les environnements simples
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Résultats")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("LES CAS ", False), ("simples", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_bullets(slide, Inches(0.55), Inches(1.75), Inches(8.9), Inches(3.2), [
    "Sur Line World, les 8 algorithmes atteignent 100% de réussite : le chemin est trop court "
    "pour qu'un algorithme se distingue d'un autre.",
    "Sur Two Round RPS, le score optimal de 1.000 est retrouvé par tous ; les quelques écarts "
    "observés viennent surtout du bruit du round 1, qui reste aléatoire.",
    "Sur Monty Hall, à 3 portes (optimal 2/3) comme à 5 portes (optimal 4/5), presque tous les "
    "algorithmes s'en approchent bien.",
    "Seule ombre au tableau : Dyna-Q est déjà le moins fiable des huit ici, avec seulement 60% "
    "de réussite sur les deux niveaux de Monty Hall.",
], size=14)
add_footer(slide)

# ---------------------------------------------------------------------------
# 10. Résultats sur Grid World
# ---------------------------------------------------------------------------
slide = new_slide(PAPER)
add_kicker(slide, "Résultats")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("GRID WORLD, LE CAS ", False), ("révélateur", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_table(
    slide, Inches(0.55), Inches(1.7), Inches(5.1), Inches(3.1),
    ["Algorithme", "Réussite"],
    [
        ["Policy / Value Iteration", "100 %"],
        ["Monte Carlo ES", "100 %"],
        ["On-policy first-visit MC", "0 %"],
        ["Off-policy MC", "100 %"],
        ["Sarsa", "80 %"],
        ["Q-Learning", "100 %"],
        ["Dyna-Q", "100 %"],
    ],
    col_widths=[Inches(3.5), Inches(1.6)],
    body_size=12,
)
add_bullets(slide, Inches(5.9), Inches(1.75), Inches(3.55), Inches(3.1), [
    "L'on-policy MC control échoue totalement : il a besoin d'un épisode complet pour "
    "apprendre quoi que ce soit, or peu d'épisodes atteignent un état terminal.",
    "Sarsa et Q-Learning apprennent à chaque pas grâce au bootstrap, ce qui les rend plus "
    "robustes face à ce même problème.",
], size=12)
add_footer(slide)

# ---------------------------------------------------------------------------
# 11. Étude des hyperparamètres
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Hyperparamètres")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("IMPACT DES ", False), ("hyperparamètres", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))

hp_blocks = [
    ("EPSILON, SUR GRID WORLD", "Plus epsilon augmente, plus Sarsa se dégrade (80% à 40% "
                                  "entre 0,1 et 0,8), alors que l'on-policy MC control fait "
                                  "l'inverse, de 0% à 80% : il a besoin de beaucoup d'exploration "
                                  "pour terminer un épisode ne serait-ce qu'une fois."),
    ("NOMBRE DE PAS DE PLANIFICATION, POUR DYNA-Q", "À nombre d'épisodes réels identique "
                                  "(500), Dyna-Q sans planification revient à du Q-Learning et "
                                  "atteint 80% de réussite. Dès 5 pas de planification, il "
                                  "atteint 100%."),
    ("GAMMA, AVEC VALUE ITERATION SUR GRID WORLD", "La valeur calculée pour l'état de "
                                  "départ passe de 0,008 avec gamma à 0,5, à 0,478 à 0,9, "
                                  "0,932 à 0,99, et proche de 1 à 0,999999."),
]
top = 1.7
for title, desc in hp_blocks:
    add_text(slide, Inches(0.55), Inches(top), Inches(8.9), Inches(0.3), title, size=13,
              bold=True, color=ORANGE, font=FONT_BODY)
    add_text(slide, Inches(0.55), Inches(top + 0.32), Inches(8.9), Inches(0.75), desc,
              size=12, color=INK, line_spacing=1.15)
    top += 1.05
add_footer(slide)

# ---------------------------------------------------------------------------
# 12. Résultats sur les environnements secrets
# ---------------------------------------------------------------------------
slide = new_slide(PAPER)
add_kicker(slide, "Résultats")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("SUR LES ENVIRONNEMENTS ", False), ("secrets", True)]], size=26, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_table(
    slide, Inches(0.55), Inches(1.7), Inches(8.9), Inches(1.75),
    ["Environnement", "Meilleur algorithme", "Score moyen"],
    [
        ["Secret Env 0 (8 192 états)", "Dyna-Q / MC on-policy", "10.0"],
        ["Secret Env 1 (65 536 états)", "Dyna-Q", "31.0"],
        ["Secret Env 2 (2 097 152 états)", "Dyna-Q", "-14.0 (budget insuffisant)"],
        ["Secret Env 3 (65 536 états)", "Sarsa", "14.8"],
    ],
    col_widths=[Inches(3.6), Inches(2.7), Inches(2.6)],
    body_size=11.5,
)
add_bullets(slide, Inches(0.55), Inches(3.7), Inches(8.9), Inches(1.3), [
    "On a trouvé et corrigé un vrai bug : evaluate_policy et replay_policy choisissaient "
    "l'action gloutonne sans vérifier qu'elle était disponible dans l'état courant, ce qui "
    "faisait planter le programme sur Secret Env 3. La correction restreint ce choix aux "
    "actions réellement jouables.",
    "Dyna-Q confirme l'avantage de la planification sur les environnements 0 et 1 ; sur "
    "l'environnement 2, le budget d'itérations reste trop faible pour un espace de 2 "
    "millions d'états.",
], size=11.5)
add_footer(slide)

# ---------------------------------------------------------------------------
# 13. Quel algorithme choisir, et pourquoi
# ---------------------------------------------------------------------------
slide = new_slide(WHITE)
add_kicker(slide, "Synthèse")
add_headline(slide, Inches(0.55), Inches(0.7), Inches(8.9), Inches(0.8),
              [[("QUEL ALGORITHME ", False), ("choisir", True), (", ET POURQUOI", False)]],
              size=25, color=INK)
add_rule(slide, Inches(0.55), Inches(1.45), Inches(8.9))
add_bullets(slide, Inches(0.55), Inches(1.7), Inches(8.9), Inches(3.3), [
    "Quand le modèle de l'environnement est connu à l'avance, Policy et Value Iteration "
    "restent le choix le plus sûr et le plus rapide, sans avoir besoin d'explorer.",
    "Quand le modèle est inconnu, Q-Learning est la méthode la plus robuste dans l'ensemble.",
    "Dyna-Q apprend plus vite que Q-Learning à nombre d'épisodes égal, mais reste le moins "
    "fiable sur Monty Hall, sans doute parce que son modèle appris suppose un environnement "
    "déterministe.",
    "L'on-policy first-visit MC control est mal adapté à Grid World, faute d'épisodes complets.",
    "Monte Carlo ES et l'off-policy MC control restent fiables partout, mais MC ES exige de "
    "pouvoir se replacer directement dans un état donné.",
    "On assume aussi nos limites : Secret Env 2, avec ses 2 millions d'états, n'est pas "
    "résolu avec le budget donné, l'interface graphique des environnements secrets n'est pas "
    "encore fournie, et Dyna-Q+ n'a pas été implémenté faute de temps.",
], size=13)
add_footer(slide)

# ---------------------------------------------------------------------------
# 14. Merci / démonstration
# ---------------------------------------------------------------------------
slide = new_slide(ORANGE)
add_kicker(slide, "Démonstration", color=WHITE)
add_headline(slide, Inches(0.55), Inches(1.8), Inches(8.9), Inches(1.6),
              [
                  [("MERCI, ", False), ("des questions ?", True)],
              ], size=38, color=WHITE)
add_text(slide, Inches(0.55), Inches(3.15), Inches(8.9), Inches(0.6),
          "Pour finir, une démonstration en direct : le rejeu pas-à-pas d'une politique "
          "déjà entraînée, sans réapprentissage.",
          size=15, color=WHITE, italic=True, font=FONT_BODY)
add_rule(slide, Inches(0.55), Inches(5.0), Inches(8.9), color=WHITE, weight=1)
add_text(slide, Inches(0.55), Inches(5.08), Inches(6.5), Inches(0.35), GROUP, size=10, color=WHITE, font=FONT_BODY)
add_text(slide, Inches(8.0), Inches(5.08), Inches(1.45), Inches(0.35), "Juillet 2026", size=10,
          color=WHITE, font=FONT_BODY, align=PP_ALIGN.RIGHT)

prs.save(OUTPUT_PATH)
print(f"Slides écrites dans {OUTPUT_PATH}")
